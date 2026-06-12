#!/usr/bin/env python3
"""Build the LocalGuard MVP proposal deck on top of the official Dacon template.

Template chrome (headers, footers, logos, section titles, divider lines) is preserved.
Only the [작성방법] guide boxes are removed; content is added inside the content zone.
Design tokens follow the Claude.com-inspired system (cream/coral/dark) with Pretendard.
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- tokens
CANVAS = RGBColor(0xFA, 0xF9, 0xF5)
CARD = RGBColor(0xEF, 0xE9, 0xDE)
CARD_SOFT = RGBColor(0xF5, 0xF0, 0xE8)
STRONG = RGBColor(0xE8, 0xE0, 0xD2)
DARK = RGBColor(0x18, 0x17, 0x15)
DARK_ELEV = RGBColor(0x25, 0x23, 0x20)
HAIR = RGBColor(0xE6, 0xDF, 0xD8)
INK = RGBColor(0x14, 0x14, 0x13)
BODY = RGBColor(0x3D, 0x3D, 0x3A)
MUTED = RGBColor(0x6C, 0x6A, 0x64)
MUTED_SOFT = RGBColor(0x8E, 0x8B, 0x82)
CORAL = RGBColor(0xCC, 0x78, 0x5C)
CORAL_DK = RGBColor(0xA9, 0x58, 0x3E)
TEAL = RGBColor(0x5D, 0xB8, 0xA6)
AMBER = RGBColor(0xE8, 0xA5, 0x5A)
ON_DARK = RGBColor(0xFA, 0xF9, 0xF5)
ON_DARK_SOFT = RGBColor(0xA0, 0x9D, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Pretendard"
FONT_SB = "Pretendard SemiBold"

IN = Inches

# ---------------------------------------------------------------- helpers

def set_run(run, text, size, color=BODY, bold=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=2):
    """lines: list of (text, size, color, bold) or list of list-of-run-tuples."""
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        runs = line if isinstance(line, list) else [line]
        for spec in runs:
            text, size, color, bold = spec[0], spec[1], spec[2], spec[3]
            font = spec[4] if len(spec) > 4 else FONT
            set_run(p.add_run(), text, size, color, bold, font)
    return box


def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.75, radius=0.08, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, weight=1.4):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1), IN(x2), IN(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    conn.shadow.inherit = False
    return conn


def add_chip(slide, x, y, w, h, text, fill, text_color, size=9, bold=True, line=None):
    shp = add_rect(slide, x, y, w, h, fill, line=line, radius=0.5)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, size, text_color, bold)
    return shp


def add_picture_card(slide, path, x, y, w, h):
    pic = slide.shapes.add_picture(path, IN(x), IN(y), IN(w), IN(h))
    pic.line.color.rgb = HAIR
    pic.line.width = Pt(1)
    pic.shadow.inherit = False
    return pic


def fill_cell(cell, text, size=14, color=BODY, bold=False, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size, color, bold)


def remove_guide_boxes(slide):
    for shp in list(slide.shapes):
        if shp.has_text_frame and "[작성방법]" in shp.text_frame.text:
            shp._element.getparent().remove(shp._element)


def style_chart(chart, font_size=9):
    chart.font.size = Pt(font_size)
    chart.font.name = FONT
    chart.font.color.rgb = BODY


# ---------------------------------------------------------------- content zone
CX, CW = 0.95, 11.60          # content x / width
CY, CH = 1.84, 4.52           # content y / height (below divider, above bottom of rect)

prs = Presentation("template.pptx")
S = list(prs.slides)

# ================================================================ Slide 1 — 표지
s = S[0]
for shp in s.shapes:
    if shp.has_table:
        tbl = shp.table
        fill_cell(tbl.cell(0, 1), "LocalGuard", 18, INK, True, PP_ALIGN.CENTER)
        fill_cell(tbl.cell(1, 1), "JB LocalGuard OS", 18, INK, True, PP_ALIGN.CENTER)

add_text(s, 1.73, 3.95, 10.4, 0.45,
         [("지역 금융 위험을 케이스로 만들고, AI Agent가 판단하고, 사람이 승인하는 금융안전 운영체제", 14, MUTED, False)],
         align=PP_ALIGN.CENTER)

# ================================================================ Slide 2 — Summary
s = S[1]
for shp in s.shapes:
    if shp.has_table:
        tbl = shp.table
        fill_cell(tbl.cell(0, 1), "LocalGuard", 15, INK, True)
        fill_cell(tbl.cell(1, 1), "이승보 (1인 팀)", 15, BODY, False)
        fill_cell(tbl.cell(2, 1), "JB LocalGuard OS  —  지역금융 안전 AI Agent 운영 콘솔", 15, BODY, True)
        fill_cell(tbl.cell(3, 1), "자유주제 (JB 미래사업 AI · 지역금융 리스크 대응 AI Agent OS)", 15, BODY, False)
        cell = tbl.cell(4, 1)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        set_run(p.add_run(), "자금압박 · 전세사기 · 보이스피싱 등 지역 고객의 금융 위험을 ", 14, BODY)
        set_run(p.add_run(), "Case", 14, CORAL_DK, True)
        set_run(p.add_run(), "로 등록하면, 전문 AI Agent들이 ", 14, BODY)
        set_run(p.add_run(), "판단 → 조치 초안 → 사람 승인 → 감사 기록", 14, CORAL_DK, True)
        set_run(p.add_run(), "까지 수행하는 JB 금융안전 운영 콘솔", 14, BODY)

# ================================================================ Slide 3 — 문제 정의
s = S[2]
remove_guide_boxes(s)

stat_specs = [
    ("8,545억 원", "2024년 보이스피싱 피해액 · 역대 최대", "전년 대비 약 2배 급증  |  경찰청 국가수사본부", CORAL_DK),
    ("39,121건", "전세사기 피해자 누적 결정 (’26.5 기준)", "’23.6 특별법 시행 후 계속 증가  |  국토교통부", CORAL_DK),
    ("43.7만 명", "취약 자영업자 (’25.2분기) · 대출 약 130조 원", "취약차주 연체율 11%대  |  한국은행 금융안정 상황", CORAL_DK),
]
sw = (CW - 0.4) / 3
for i, (num, label, src, accent) in enumerate(stat_specs):
    x = CX + i * (sw + 0.2)
    add_rect(s, x, CY, sw, 1.12, CARD_SOFT, line=HAIR)
    add_text(s, x + 0.18, CY + 0.10, sw - 0.36, 0.42, [(num, 22, accent, True)])
    add_text(s, x + 0.18, CY + 0.52, sw - 0.36, 0.30, [(label, 10.5, INK, True)])
    add_text(s, x + 0.18, CY + 0.82, sw - 0.36, 0.24, [(src, 8, MUTED_SOFT, False)])

# charts row
chy, chh = CY + 1.28, 2.08
chw = (CW - 0.3) / 2

add_rect(s, CX, chy, chw, chh, WHITE, line=HAIR)
add_text(s, CX + 0.18, chy + 0.08, chw - 0.36, 0.25,
         [("보이스피싱 피해액 추이 (억 원)", 11, INK, True)])
cd = CategoryChartData()
cd.categories = ["2021", "2022", "2023", "2024"]
cd.add_series("피해액(억원)", (7744, 5438, 4472, 8545))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        IN(CX + 0.15), IN(chy + 0.36), IN(chw - 0.3), IN(chh - 0.72), cd)
chart = gf.chart
chart.has_legend = False
style_chart(chart)
plot = chart.plots[0]
plot.gap_width = 80
ser = plot.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = CORAL
pts = ser.points
for idx in (0, 1, 2):
    pts[idx].format.fill.solid()
    pts[idx].format.fill.fore_color.rgb = STRONG
plot.has_data_labels = True
plot.data_labels.font.size = Pt(9)
plot.data_labels.font.color.rgb = INK
plot.data_labels.font.name = FONT
chart.value_axis.visible = False
chart.value_axis.has_major_gridlines = False
chart.category_axis.format.line.color.rgb = HAIR
add_text(s, CX + 0.18, chy + chh - 0.30, chw - 0.36, 0.22,
         [("출처: 경찰청 국가수사본부 발표 (2025.7 보도) · 2024년 1인당 피해액 5,290만 원으로 2배 증가", 8, MUTED_SOFT, False)])

x2 = CX + chw + 0.3
add_rect(s, x2, chy, chw, chh, WHITE, line=HAIR)
add_text(s, x2 + 0.18, chy + 0.08, chw - 0.36, 0.25,
         [("전세사기 피해자 누적 결정 건수", 11, INK, True)])
cd2 = CategoryChartData()
cd2.categories = ["’25.8", "’25.10", "’25.11", "’26.5"]
cd2.add_series("누적 결정(건)", (32185, 33978, 34481, 39121))
gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                         IN(x2 + 0.15), IN(chy + 0.36), IN(chw - 0.3), IN(chh - 0.72), cd2)
chart2 = gf2.chart
chart2.has_legend = False
style_chart(chart2)
ser2 = chart2.plots[0].series[0]
ser2.format.line.color.rgb = CORAL_DK
ser2.format.line.width = Pt(2.25)
chart2.plots[0].has_data_labels = True
chart2.plots[0].data_labels.font.size = Pt(8.5)
chart2.plots[0].data_labels.font.color.rgb = INK
chart2.plots[0].data_labels.font.name = FONT
chart2.value_axis.visible = False
chart2.value_axis.has_major_gridlines = False
chart2.category_axis.format.line.color.rgb = HAIR
add_text(s, x2 + 0.18, chy + chh - 0.30, chw - 0.36, 0.22,
         [("출처: 국토교통부 전세사기피해지원위원회 결정 현황 (보도자료 ’25.8 / ’25.10 / ’25.11 / ’26.5)", 8, MUTED_SOFT, False)])

# bottom problem strip
by = chy + chh + 0.16
add_rect(s, CX, by, CW, 1.0, DARK, radius=0.10)
cols = [
    ("누가 겪는 문제인가", "지역 소상공인 · 전세 계약을 앞둔 청년 · 고령 고객,\n그리고 이들을 담당하는 JB RM · 심사 · 준법 담당자"),
    ("무엇이 문제인가", "기사 · 사기 경보 · 시세 · 등기 · 상담기록이 흩어져 있어\n위험 신호를 조기에 모아서 판단하고 행동하기 어려움"),
    ("해결 후 기대 변화 (KPI)", "위험 인지→대응 착수 시간 50% 단축 · 판단 100% 근거 연결\n고객 대상 행동 100% 사람 승인 통과 · 사후관리 누락 0건"),
]
cw3 = (CW - 0.9) / 3
for i, (t, b) in enumerate(cols):
    x = CX + 0.25 + i * (cw3 + 0.2)
    add_text(s, x, by + 0.13, cw3, 0.26, [(t, 10.5, CORAL, True)])
    add_text(s, x, by + 0.40, cw3, 0.55, [(line, 9, ON_DARK_SOFT, False) for line in b.split("\n")], space_after=1)

# ================================================================ Slide 4 — 솔루션 개요
s = S[3]
remove_guide_boxes(s)

# left: architecture diagram
ax, aw = CX, 7.1
add_text(s, ax, CY - 0.02, aw, 0.26, [("시스템 구성도 — Case → AgentRun → Approval → Audit 운영 루프", 11.5, INK, True)])

# L1 input layer
l1y = CY + 0.30
add_rect(s, ax, l1y, aw, 0.62, CARD_SOFT, line=HAIR)
add_text(s, ax + 0.12, l1y + 0.06, 1.1, 0.5, [("입력 · 증거", 9.5, MUTED, True)], anchor=MSO_ANCHOR.MIDDLE)
in_chips = ["상담 노트", "뉴스 · 기사", "공식 경보(금융위·국토부)", "시세 · 등기(HUG 안심전세)"]
icw = [1.05, 1.05, 1.85, 1.75]
ix = ax + 1.25
for label, w in zip(in_chips, icw):
    add_chip(s, ix, l1y + 0.14, w, 0.34, label, WHITE, BODY, size=8.5, bold=False, line=HAIR)
    ix += w + 0.1

# arrow L1 -> L2
add_arrow(s, ax + aw / 2, l1y + 0.62, ax + aw / 2, l1y + 0.82)

# L2 dark core
l2y = l1y + 0.84
l2h = 1.74
add_rect(s, ax, l2y, aw, l2h, DARK, radius=0.06)
add_text(s, ax + 0.16, l2y + 0.09, 4.4, 0.24, [("JB LocalGuard OS  ·  멀티 Agent 코어", 10.5, ON_DARK, True)])
add_text(s, ax + 4.6, l2y + 0.10, aw - 4.75, 0.22, [("Skill Registry 25종 장착", 8.5, AMBER, True)], align=PP_ALIGN.RIGHT)
add_chip(s, ax + 0.16, l2y + 0.40, 2.30, 0.34, "LocalGuard Orchestrator", DARK_ELEV, ON_DARK, size=9)
add_text(s, ax + 2.52, l2y + 0.44, aw - 2.7, 0.24, [("지시 해석 · Agent 배정 · 승인 레벨(L0–L4) 산정", 8.5, ON_DARK_SOFT, False)])
agent_chips = ["Pain Radar", "Cashflow Triage", "Policy Match", "Fraud Shield", "RM Copilot", "Compliance Guard"]
acw = (aw - 0.32 - 0.5) / 6
for i, label in enumerate(agent_chips):
    add_chip(s, ax + 0.16 + i * (acw + 0.1), l2y + 0.84, acw, 0.34, label, DARK_ELEV, ON_DARK, size=8)
add_rect(s, ax + 0.16, l2y + 1.26, aw - 0.32, 0.36, DARK_ELEV, radius=0.5)
add_text(s, ax + 0.32, l2y + 1.33, aw - 0.6, 0.24,
         [[("Jeonse Shield 라인 (5 Agent)  ", 8.5, TEAL, True),
           ("전세가율 · 권리관계 · 자산 리스크 · 계약 체크리스트 · 은행 연계", 8.5, ON_DARK_SOFT, False)]])

# arrow L2 -> L3
add_arrow(s, ax + aw / 2, l2y + l2h, ax + aw / 2, l2y + l2h + 0.20)

# L3 approval gate (coral)
l3y = l2y + l2h + 0.22
add_rect(s, ax, l3y, aw, 0.52, CORAL, radius=0.08)
add_text(s, ax + 0.16, l3y + 0.13, aw - 0.32, 0.28,
         [[("Approval Gate  ", 11, WHITE, True),
           ("고객 대상 행동 · 민감 조치는 사람(RM · 준법)이 승인해야 실행", 9.5, WHITE, False)]])

# arrow L3 -> L4
add_arrow(s, ax + aw / 2, l3y + 0.52, ax + aw / 2, l3y + 0.72)

# L4 output
l4y = l3y + 0.74
half = (aw - 0.2) / 2
add_rect(s, ax, l4y, half, 0.54, CARD, line=HAIR)
add_text(s, ax + 0.14, l4y + 0.07, half - 0.28, 0.4,
         [("Audit Ledger", 9.5, INK, True), ("근거 · 판단 · 행동 · 승인 전체 기록", 8, MUTED, False)], space_after=0)
add_rect(s, ax + half + 0.2, l4y, half, 0.54, CARD, line=HAIR)
add_text(s, ax + half + 0.34, l4y + 0.07, half - 0.28, 0.4,
         [("RM 콘솔 · 고객 안내(승인 후)", 9.5, INK, True), ("콜백 초안 · 체크리스트 · 은행 서비스 연계", 8, MUTED, False)], space_after=0)

# right column: operating loop + principles
rx = CX + aw + 0.35
rw = CW - aw - 0.35
add_text(s, rx, CY - 0.02, rw, 0.26, [("서비스 구성 · 운영 시나리오", 11.5, INK, True)])
loop_steps = [
    ("1", "Case 생성", "고객 위험을 케이스로 등록 (RM 또는 정기 스캔)"),
    ("2", "AgentRun", "전문 Agent들이 스킬을 장착하고 근거 수집 · 위험 판단"),
    ("3", "Approval", "조치 초안을 사람이 검토 — Approve / Reject"),
    ("4", "Audit", "모든 판단 · 행동 · 승인이 감사 로그로 남음"),
]
ly = CY + 0.30
for num, t, b in loop_steps:
    add_rect(s, rx, ly, rw, 0.56, CANVAS, line=HAIR)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(rx + 0.10), IN(ly + 0.13), IN(0.30), IN(0.30))
    circ.fill.solid()
    circ.fill.fore_color.rgb = CORAL
    circ.line.fill.background()
    circ.shadow.inherit = False
    tfc = circ.text_frame
    tfc.margin_left = tfc.margin_right = tfc.margin_top = tfc.margin_bottom = 0
    pc = tfc.paragraphs[0]
    pc.alignment = PP_ALIGN.CENTER
    set_run(pc.add_run(), num, 10, WHITE, True)
    add_text(s, rx + 0.52, ly + 0.06, rw - 0.65, 0.22, [(t, 10, INK, True)])
    add_text(s, rx + 0.52, ly + 0.28, rw - 0.65, 0.22, [(b, 8.5, MUTED, False)])
    ly += 0.64

add_rect(s, rx, ly + 0.02, rw, 1.62, DARK, radius=0.08)
add_text(s, rx + 0.16, ly + 0.13, rw - 0.32, 0.24, [("핵심 설계 원칙", 10, CORAL, True)])
principles = [
    ("스킬 장착형 멀티 Agent", "기능 추가 = 스킬 패키지 추가 (확장 용이)"),
    ("승인 우선 자동화 (L0–L4)", "완전 자동 발송 금지 · 사람이 마지막 결정"),
    ("근거 · 감사가 일급 객체", "모든 판단에 출처 링크와 감사 기록 연결"),
]
py = ly + 0.40
for t, b in principles:
    add_text(s, rx + 0.16, py, rw - 0.32, 0.22, [[("· " + t, 9, ON_DARK, True)]])
    add_text(s, rx + 0.30, py + 0.20, rw - 0.46, 0.20, [(b, 8, ON_DARK_SOFT, False)])
    py += 0.41

# ================================================================ Slide 5 — 주요 기능
s = S[4]
remove_guide_boxes(s)

shot_w = 5.55
shot_h = shot_w * 9 / 16
add_text(s, CX, CY - 0.02, shot_w, 0.26, [("실제 동작하는 MVP 콘솔 (브라우저 데모)", 11.5, INK, True)])
add_picture_card(s, "shots/dashboard.png", CX, CY + 0.30, shot_w, shot_h)
add_text(s, CX, CY + 0.34 + shot_h, shot_w, 0.22,
         [("대시보드 — 지시 입력(Dispatch) · 라이브 AgentRun · 위험 지표 · 승인 큐", 8.5, MUTED, False)])
strip_y = CY + 0.64 + shot_h
add_rect(s, CX, strip_y, shot_w, 0.74, DARK, radius=0.10)
add_text(s, CX + 0.18, strip_y + 0.10, shot_w - 0.36, 0.24,
         [[("MVP 범위: ", 9.5, CORAL, True),
           ("판단(위험 진단) → 행동(조치 초안) → 검증(승인·감사)", 9.5, ON_DARK, True)]])
add_text(s, CX + 0.18, strip_y + 0.40, shot_w - 0.36, 0.24,
         [("챌린지가 요구하는 AI Agent 3요소를 한 화면에서 시연 — github.com/LSB-afk/JB-Fin-AI-Challenge", 8, ON_DARK_SOFT, False)])

fx = CX + shot_w + 0.35
fw = CW - shot_w - 0.35
add_text(s, fx, CY - 0.02, fw, 0.26, [("핵심 기능 6가지 — 모두 MVP에서 동작 확인", 11.5, INK, True)])
features = [
    ("Case 보드", "위험 케이스를 칸반/큐로 관리, 클릭 시 근거·담당 Agent 연동", "흩어진 고객 위험의 단일 작업 화면"),
    ("AgentRun 실행 로그", "Run/Dispatch마다 실행 로그 생성, 상태 전이 자동 기록", "AI 판단 과정의 투명한 추적"),
    ("승인 큐", "Approve/Reject 클릭으로 상태·감사·큐 숫자 즉시 갱신", "고객 대상 행동의 사람 통제"),
    ("전세 Shield", "전세가율·권리관계·자산리스크·체크리스트·은행연계 5대 진단", "전세사기 계약 전 예방"),
    ("Skill Registry", "25종 스킬에 위험도·승인 정책 부여, Agent에 장착", "기능 확장과 내부통제의 단위화"),
    ("Audit Ledger", "근거→판단→행동→승인 전 과정을 시간순 기록", "금융권 감사·설명가능성 충족"),
]
fcw = (fw - 0.2) / 2
fch = 1.26
for i, (t, d, p) in enumerate(features):
    col, row = i % 2, i // 2
    x = fx + col * (fcw + 0.2)
    y = CY + 0.30 + row * (fch + 0.12)
    add_rect(s, x, y, fcw, fch, CARD_SOFT, line=HAIR)
    add_chip(s, x + fcw - 1.18, y + 0.09, 1.06, 0.24, "MVP 구현 완료", TEAL, WHITE, size=7.5)
    add_text(s, x + 0.13, y + 0.08, fcw - 1.35, 0.24, [(t, 10.5, INK, True)])
    add_text(s, x + 0.13, y + 0.36, fcw - 0.26, 0.50, [(d, 8.5, BODY, False)])
    add_text(s, x + 0.13, y + 0.94, fcw - 0.26, 0.26, [[("해결: ", 8, CORAL_DK, True), (p, 8, MUTED, False)]])

# ================================================================ Slide 6 — 데이터 및 기술
s = S[5]
remove_guide_boxes(s)

dw = 6.0
add_text(s, CX, CY - 0.02, dw, 0.26, [("활용 데이터 — 공개 · 공식 출처 우선, 출처 전부 명시", 11.5, INK, True)])
data_rows = [
    ("공공 · 공식", "HUG 안심전세(시세·전세가율·보증) · 국토교통부 보도자료 · 금융위 보이스피싱 경보", "공개 자료 · 출처 표기"),
    ("시장 · 언론", "연합뉴스 · 쿠키뉴스 등 지역 경기/금리 부담 기사 (Pain Radar 수집 대상)", "링크 단위 근거 연결"),
    ("JB 공식", "JB금융그룹 계열사 현황 · JB-네이버클라우드 AI 업무협약 (jbfg.com)", "사업 연계성 근거"),
    ("내부 운영(목업)", "상담 노트 · 케이스 이력 · 승인 기록 — MVP는 가상 데이터로 구성", "개인정보 마스킹 전제"),
]
dy = CY + 0.30
for t, b, m in data_rows:
    add_rect(s, CX, dy, dw, 0.70, CANVAS, line=HAIR)
    add_text(s, CX + 0.14, dy + 0.08, 1.35, 0.5, [(t, 9.5, CORAL_DK, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, CX + 1.55, dy + 0.08, dw - 3.0, 0.56, [(b, 8.5, BODY, False)])
    add_text(s, CX + dw - 1.42, dy + 0.08, 1.30, 0.56, [(m, 7.5, MUTED_SOFT, False)], align=PP_ALIGN.RIGHT)
    dy += 0.78

tx = CX + dw + 0.35
tw = CW - dw - 0.35
add_text(s, tx, CY - 0.02, tw, 0.26, [("기술 구성 — 현재 MVP와 본선 목표", 11.5, INK, True)])
add_rect(s, tx, CY + 0.30, tw, 2.96, DARK, radius=0.06)
add_text(s, tx + 0.18, CY + 0.42, tw - 0.36, 0.22, [("$ 현재 MVP  (예선 · 구현 완료)", 9.5, TEAL, True, "Pretendard")])
mvp_lines = [
    "Vanilla JS/CSS 정적 콘솔 — 외부 의존성 0, 브라우저만으로 시연",
    "상태 모델: Case · AgentRun · Approval · Audit 상호 연동",
    "GitHub 공개: github.com/LSB-afk/JB-Fin-AI-Challenge",
]
my = CY + 0.66
for line in mvp_lines:
    add_text(s, tx + 0.30, my, tw - 0.48, 0.2, [("· " + line, 8.5, ON_DARK_SOFT, False)])
    my += 0.24
add_text(s, tx + 0.18, my + 0.10, tw - 0.36, 0.22, [("$ 본선 목표 아키텍처", 9.5, AMBER, True)])
tech_lines = [
    "LLM — 위험 판단 요약 · 조치 초안 생성 (Claude 등 상용 API)",
    "RAG — HUG·국토부·금융위 공개자료 검색으로 근거 강제 연결",
    "Rule Engine — 승인 레벨(L0–L4) · 금지 행동 정책 판정",
    "Multi-Agent — Orchestrator가 전문 Agent에 스킬 주입 · 배정",
    "외부 API — 공공 데이터 우선, 상용 API는 라이선스 검토 후 사용",
]
ty2 = my + 0.36
for line in tech_lines:
    add_text(s, tx + 0.30, ty2, tw - 0.48, 0.2, [("· " + line, 8.5, ON_DARK_SOFT, False)])
    ty2 += 0.24

ry = CY + 3.44
add_text(s, CX, ry, CW, 0.24, [("기술적 제약과 해결 전략", 11.5, INK, True)])
limits = [
    ("환각(Hallucination)", "모든 판단에 출처 링크 강제 + 확정 표현 금지 + 사람 승인 게이트", CORAL_DK),
    ("개인정보 · 내부통제", "마스킹 처리 · L0–L4 승인 레벨 · 전 과정 감사 로그", CORAL_DK),
    ("법률 판단 한계", "등기·보증은 원문 확인 절차 명시, 특약은 초안만 제공(법무 검토)", CORAL_DK),
]
lw3 = (CW - 0.4) / 3
for i, (t, b, c) in enumerate(limits):
    x = CX + i * (lw3 + 0.2)
    add_rect(s, x, ry + 0.28, lw3, 0.72, CARD, line=HAIR)
    add_text(s, x + 0.14, ry + 0.36, lw3 - 0.28, 0.22, [(t, 9.5, c, True)])
    add_text(s, x + 0.14, ry + 0.58, lw3 - 0.28, 0.4, [(b, 8.5, BODY, False)])

# ================================================================ Slide 7 — 사용자 시나리오
s = S[6]
remove_guide_boxes(s)

add_text(s, CX, CY - 0.02, CW, 0.26,
         [[("대표 유스케이스 — 전세 Shield  ", 11.5, INK, True),
           ("(청년 고객이 전세계약 전에 은행을 통해 위험을 진단받는 흐름)", 9.5, MUTED, False)]])

actors = [("고객", CARD, INK), ("RM(은행)", CORAL, WHITE), ("AI Agent", DARK, ON_DARK), ("검증·승인", TEAL, WHITE)]
lx = CX
for label, fillc, tc in actors:
    add_chip(s, lx, CY + 0.28, 0.95, 0.26, label, fillc, tc, size=8)
    lx += 1.05

steps = [
    ("①", "전세계약 전\n상담 신청", "고객", CARD, INK),
    ("②", "Case 생성\nJBG-201", "RM(은행)", CORAL, WHITE),
    ("③", "전세가율·권리관계\n·자산 리스크 진단", "AI Agent", DARK, ON_DARK),
    ("④", "체크리스트 +\n특약 초안 생성", "AI Agent", DARK, ON_DARK),
    ("⑤", "Compliance\n표현·준법 검토", "검증·승인", TEAL, WHITE),
    ("⑥", "RM 승인\nApprove", "RM(은행)", CORAL, WHITE),
    ("⑦", "안전 계약 가이드\n+ 대출·보증 상담", "고객", CARD, INK),
]
fy = CY + 0.66
fh = 1.06
gap = 0.24
fw7 = (CW - gap * 6) / 7
for i, (num, label, actor, fillc, tc) in enumerate(steps):
    x = CX + i * (fw7 + gap)
    add_rect(s, x, fy, fw7, fh, fillc, line=HAIR if fillc in (CARD,) else None, radius=0.12)
    add_text(s, x + 0.08, fy + 0.08, fw7 - 0.16, 0.24, [(num, 11, tc, True)])
    add_text(s, x + 0.08, fy + 0.34, fw7 - 0.16, 0.6,
             [(seg, 8.5, tc, True) for seg in label.split("\n")], space_after=0)
    if i < 6:
        add_arrow(s, x + fw7 + 0.015, fy + fh / 2, x + fw7 + gap - 0.015, fy + fh / 2, MUTED, 1.2)

add_text(s, CX, fy + fh + 0.06, CW, 0.22,
         [[("안전 정책: ", 8.5, CORAL_DK, True),
           ("법률 확정 표현 금지 · 등기/보증은 원문 확인 · 특약은 초안만 · 은행 연계는 고객 동의 + RM 승인 후 진행", 8.5, MUTED, False)]])

sy = fy + fh + 0.40
shot2_w = 5.55
shot2_h = shot2_w * 9 / 16 * 0.71
add_picture_card(s, "shots/jeonse.png", CX, sy, shot2_w, shot2_h)
add_text(s, CX, sy + shot2_h + 0.04, shot2_w, 0.2,
         [("MVP 전세 Shield 화면 — 5대 진단 기능과 케이스(JBG-201) 연동", 8, MUTED, False)])

bx = CX + shot2_w + 0.35
bw = CW - shot2_w - 0.35
add_rect(s, bx, sy, bw, 2.38, CARD_SOFT, line=HAIR)
add_text(s, bx + 0.16, sy + 0.10, bw - 0.32, 0.24, [("보조 유스케이스 — 소상공인 자금압박 케어", 10, INK, True)])
sub_steps = [
    "전주 카페 사장님 케이스에 매출 둔화 + 금리 부담 신호 감지 (Pain Radar)",
    "Cashflow Triage가 상환 스트레스 판단 → Policy Match가 정책금융·서류 체크리스트 생성",
    "RM Copilot이 콜백 초안 작성 → Compliance 검토 → RM 승인 후 고객 안내",
    "전 과정이 Audit Ledger에 기록 — 사후관리 누락과 책임소재 문제 해소",
]
sy2 = sy + 0.40
for i, line in enumerate(sub_steps, 1):
    add_text(s, bx + 0.16, sy2, bw - 0.32, 0.34,
             [[(f"{i}. ", 8.5, CORAL_DK, True), (line, 8.5, BODY, False)]], space_after=0)
    sy2 += 0.46

# ================================================================ Slide 8 — 기대 효과
s = S[7]
remove_guide_boxes(s)

impact = [
    ("50%↓", "위험 인지→대응 착수 시간", "RM이 케이스 화면 하나로 원인·근거·다음 행동 파악"),
    ("100%", "판단의 근거 연결률", "모든 Agent 판단에 출처 링크 또는 내부 이벤트 연결"),
    ("100%", "고객 대상 행동 승인 통과율", "승인 게이트 없이는 외부 안내 불가 — 내부통제 내장"),
    ("0건", "고위험 사기 케이스 외부 발송", "Fraud Shield가 자동 차단, 보안 escalation만 허용"),
]
iw = (CW - 0.6) / 4
for i, (num, t, b) in enumerate(impact):
    x = CX + i * (iw + 0.2)
    add_rect(s, x, CY, iw, 1.28, CARD_SOFT, line=HAIR)
    add_text(s, x + 0.14, CY + 0.10, iw - 0.28, 0.36, [(num, 20, CORAL_DK, True)])
    add_text(s, x + 0.14, CY + 0.50, iw - 0.28, 0.24, [(t, 9, INK, True)])
    add_text(s, x + 0.14, CY + 0.78, iw - 0.28, 0.44, [(b, 8, MUTED, False)])
add_text(s, CX, CY + 1.32, CW, 0.2,
         [("* MVP 운영 목표 지표 — 본선에서 실데이터 PoC로 측정·검증 예정", 8, MUTED_SOFT, False)])

ry2 = CY + 1.66
rw2 = 6.6
add_text(s, CX, ry2, rw2, 0.24, [("발전 경로 — 예선에서 고객 서비스까지", 11.5, INK, True)])
road = [
    ("예선 (완료)", "정적 MVP 콘솔\n7개 시나리오 검증 · GitHub 공개", TEAL),
    ("본선", "LLM · RAG 연동\n공공데이터 PoC", CORAL),
    ("내부 파일럿", "전북은행 RM 부서\n승인·감사 정책 검증", CORAL_DK),
    ("고객 서비스화", "뱅킹 앱 연계\n전세 Shield 우선 출시", DARK),
]
rcw = (rw2 - 0.45) / 4
for i, (t, b, c) in enumerate(road):
    x = CX + i * (rcw + 0.15)
    add_rect(s, x, ry2 + 0.30, rcw, 1.12, WHITE, line=HAIR)
    add_rect(s, x, ry2 + 0.30, rcw, 0.30, c, radius=0.15)
    add_text(s, x + 0.08, ry2 + 0.345, rcw - 0.16, 0.22, [(t, 8.5, WHITE if c != CARD else INK, True)], align=PP_ALIGN.CENTER)
    add_text(s, x + 0.10, ry2 + 0.70, rcw - 0.20, 0.66,
             [(seg, 8, BODY, False) for seg in b.split("\n")], space_after=1)
    if i < 3:
        add_arrow(s, x + rcw + 0.005, ry2 + 0.86, x + rcw + 0.145, ry2 + 0.86, MUTED, 1.2)

kx = CX + rw2 + 0.35
kw = CW - rw2 - 0.35
add_text(s, kx, ry2, kw, 0.24, [("운영 리스크 대응", 11.5, INK, True)])
risks = [
    ("환각 · 오판", "근거 링크 강제 + 확정 표현 금지 + 사람 승인"),
    ("개인정보 · 보안", "마스킹 · 내부망 전제 · 접근 권한 분리"),
    ("책임소재 · 설명가능성", "감사 로그로 판단·승인 주체 전 과정 추적"),
    ("저작권 · 라이선스", "공공·공식 출처 우선, 기사 등은 링크 인용"),
]
ky = ry2 + 0.30
for t, b in risks:
    add_text(s, kx, ky, kw, 0.3,
             [[("· " + t + " — ", 8.7, INK, True), (b, 8.7, BODY, False)]], space_after=0)
    ky += 0.30

by8 = ry2 + 1.62
add_rect(s, CX, by8, CW, 0.78, DARK, radius=0.10)
add_text(s, CX + 0.25, by8 + 0.12, CW - 0.5, 0.26,
         [[("계열사 확장: ", 9.5, CORAL, True),
           ("전북은행(지역 SME · 전세 Shield)  →  광주은행(가계 · 청년)  →  JB우리캐피탈(할부 · 리스 사후관리)", 9.5, ON_DARK, False)]])
add_text(s, CX + 0.25, by8 + 0.42, CW - 0.5, 0.24,
         [("JB금융그룹이 공개한 AI 추진 방향(기업대출 상담·심사·사후관리 AI — 네이버클라우드 업무협약)과 같은 축에서 확장", 8.5, ON_DARK_SOFT, False)])

prs.save("[LocalGuard] JB금융그룹 Fin AI Challenge MVP제안서.pptx")
print("saved OK")
