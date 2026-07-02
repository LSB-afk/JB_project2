#!/usr/bin/env python3
"""Build the LocalGuard MVP proposal deck on top of the official Dacon template.

Template chrome (headers, footers, logos, section titles, divider lines) is preserved.
Only the [작성방법] guide boxes are removed; content is added inside the content zone.

Design system: matches the JB LocalGuard OS app (02_제품/app/styles.css) — navy/blue/cyan
on light surfaces, with deep navy as the dark product surface. Font fixed to Pretendard.
Architecture and PII-governance diagrams are drawn as NATIVE, editable pptx shapes.
Screenshots are captured fresh from the latest app (capture_shots.js).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- tokens (02_제품/app/styles.css)
INK = RGBColor(0x0F, 0x1F, 0x3A)           # --ink
BODY = RGBColor(0x24, 0x39, 0x5F)          # readable mid
MUTED = RGBColor(0x63, 0x77, 0x97)         # --muted
MUTED_SOFT = RGBColor(0x8A, 0x9A, 0xB4)
LINE = RGBColor(0xD6, 0xE3, 0xF1)          # --line
WHITE = RGBColor(0xFF, 0xFF, 0xFF)         # --surface
SOFT = RGBColor(0xF6, 0xF9, 0xFD)          # --surface-soft
TINT = RGBColor(0xED, 0xF7, 0xFF)          # --surface-blue
TINT_STRONG = RGBColor(0xCF, 0xE6, 0xFA)   # stronger blue tint (chart non-highlight)
BLUE = RGBColor(0x0B, 0x6F, 0xB3)          # --blue-600 (primary accent)
BLUE7 = RGBColor(0x07, 0x56, 0x9D)         # --blue-700 (darker accent / numbers)
NAVY8 = RGBColor(0x12, 0x3D, 0x82)         # --navy-800
DARK = RGBColor(0x09, 0x1C, 0x4D)          # --navy-950 (dark product surface)
ON_DARK = RGBColor(0xED, 0xF4, 0xFB)       # --bg, used as on-dark text
ON_DARK_SOFT = RGBColor(0x9D, 0xB1, 0xCE)
CYAN = RGBColor(0x1C, 0xA9, 0xD6)          # --cyan-500 (accent on dark)
CYAN1 = RGBColor(0xDF, 0xF7, 0xFF)         # --cyan-100 (soft tint)
SUCCESS = RGBColor(0x0F, 0x8F, 0x72)
SUCCESS_BG = RGBColor(0xE2, 0xF7, 0xF1)

FONT = "Pretendard"
IN = Inches

# ---------------------------------------------------------------- helpers

def set_run(run, text, size, color=INK, bold=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=2):
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        runs = line if isinstance(line, list) else [line]
        for spec in runs:
            text, size, color, bold = spec[0], spec[1], spec[2], spec[3]
            font = spec[4] if len(spec) > 4 else FONT
            set_run(p.add_run(), text, size, color, bold, font)
    return box


def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.75, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, weight=1.3):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1), IN(x2), IN(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    conn.shadow.inherit = False
    return conn


def add_chip(slide, x, y, w, h, text, fill, text_color, size=8.5, bold=True, line=None):
    shp = add_rect(slide, x, y, w, h, fill, line=line, radius=0.5)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, size, text_color, bold)
    return shp


def add_node(slide, x, y, w, h, title, sub, fill, title_color, sub_color,
             line=None, radius=0.12, title_size=9, sub_size=7):
    """A flow-diagram node: bold title + small caption, centered."""
    add_rect(slide, x, y, w, h, fill, line=line, radius=radius)
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(1)
    set_run(p.add_run(), title, title_size, title_color, True)
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2.add_run(), sub, sub_size, sub_color, False)
    return box


def add_picture_card(slide, path, x, y, w, h):
    pic = slide.shapes.add_picture(path, IN(x), IN(y), IN(w), IN(h))
    pic.line.color.rgb = LINE
    pic.line.width = Pt(1)
    pic.shadow.inherit = False
    return pic


def fill_cell(cell, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size, color, bold)


def remove_guide_boxes(slide):
    for shp in list(slide.shapes):
        if shp.has_text_frame and "[작성방법]" in shp.text_frame.text:
            shp._element.getparent().remove(shp._element)


def style_chart(chart, font_size=9):
    chart.font.size = Pt(font_size)
    chart.font.name = FONT
    chart.font.color.rgb = MUTED


# ---------------------------------------------------------------- content zone
CX, CW = 0.95, 11.60
CY, CH = 1.84, 4.52

prs = Presentation("template.pptx")
S = list(prs.slides)

# ================================================================ Slide 1 — 표지
s = S[0]
for shp in s.shapes:
    if shp.has_table:
        tbl = shp.table
        fill_cell(tbl.cell(0, 1), "LocalGuard", 18, INK, True, PP_ALIGN.CENTER)
        fill_cell(tbl.cell(1, 1), "JB LocalGuard OS", 18, BLUE7, True, PP_ALIGN.CENTER)

add_text(s, 1.73, 3.95, 10.4, 0.45,
         [("지역 금융 위험을 케이스로 모으고, AI 에이전트가 판단하고, 사람이 승인하는 금융안전 운영체제",
           14, MUTED, False)],
         align=PP_ALIGN.CENTER)

# ================================================================ Slide 2 — Summary
s = S[1]
for shp in s.shapes:
    if shp.has_table:
        tbl = shp.table
        fill_cell(tbl.cell(0, 1), "LocalGuard", 15, INK, True)
        fill_cell(tbl.cell(1, 1), "이승보 (1인 팀)", 15, INK, False)
        fill_cell(tbl.cell(2, 1), "JB LocalGuard OS  —  지역금융 안전 AI 에이전트 운영 콘솔", 15, INK, True)
        fill_cell(tbl.cell(3, 1), "자유주제 (JB 미래사업 AI · 지역금융 리스크 대응 AI 에이전트 OS)", 15, INK, False)
        cell = tbl.cell(4, 1)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        set_run(p.add_run(), "자금압박 · 전세사기 · 보이스피싱 등 지역 고객의 금융 위험을 ", 14, INK)
        set_run(p.add_run(), "케이스", 14, BLUE7, True)
        set_run(p.add_run(), "로 등록하면, 전문 AI 에이전트들이 ", 14, INK)
        set_run(p.add_run(), "판단 → 조치 초안 → 사람 승인 → 감사 기록", 14, BLUE7, True)
        set_run(p.add_run(), "까지 수행하는 JB 금융안전 운영 콘솔", 14, INK)

# ================================================================ Slide 3 — 문제 정의
s = S[2]
remove_guide_boxes(s)

stat_specs = [
    ("1.13조 원", "보이스피싱 피해액 (’25.1~11월)", "전년 동기 대비 +56%  ·  금융위 · 경찰청"),
    ("39,121건", "전세사기 피해자 누적 결정 (’26 기준)", "HUG 사고액 ’24년 4.49조 원(역대 최대)  ·  국토부 · HUG"),
    ("42.7만 명", "취약 자영업자 (’24말 · 13.7%)", "취약차주 연체율 11.16%  ·  한국은행 금융안정상황"),
]
sw = (CW - 0.4) / 3
for i, (num, label, src) in enumerate(stat_specs):
    x = CX + i * (sw + 0.2)
    add_rect(s, x, CY, sw, 1.00, TINT, line=LINE)
    add_rect(s, x, CY, 0.07, 1.00, BLUE, radius=0.5)
    add_text(s, x + 0.20, CY + 0.09, sw - 0.36, 0.34, [(num, 18, BLUE7, True)])
    add_text(s, x + 0.20, CY + 0.45, sw - 0.36, 0.26, [(label, 9.5, INK, True)])
    add_text(s, x + 0.20, CY + 0.71, sw - 0.36, 0.22, [(src, 7.5, MUTED, False)])

chy, chh = CY + 1.14, 2.00
chw = (CW - 0.3) / 2

add_rect(s, CX, chy, chw, chh, WHITE, line=LINE)
add_text(s, CX + 0.16, chy + 0.08, chw - 0.32, 0.22,
         [("보이스피싱 피해액 비교 (억 원, 1~11월)", 10, INK, True)])
cd = CategoryChartData()
cd.categories = ["2024년 1~11월", "2025년 1~11월"]
cd.add_series("피해액(억원)", (7263, 11330))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        IN(CX + 0.14), IN(chy + 0.34), IN(chw - 0.28), IN(chh - 0.68), cd)
chart = gf.chart
chart.has_legend = False
style_chart(chart)
plot = chart.plots[0]
plot.gap_width = 90
ser = plot.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = BLUE
pts = ser.points
pts[0].format.fill.solid()
pts[0].format.fill.fore_color.rgb = TINT_STRONG
plot.has_data_labels = True
plot.data_labels.font.size = Pt(8.5)
plot.data_labels.font.color.rgb = INK
plot.data_labels.font.name = FONT
chart.value_axis.visible = False
chart.value_axis.has_major_gridlines = False
chart.category_axis.format.line.color.rgb = LINE
add_text(s, CX + 0.16, chy + chh - 0.28, chw - 0.32, 0.2,
         [("출처: 경찰청 국가수사본부 발표 (2025.7 보도) · 2024년 1인당 피해액 5,290만 원으로 2배 증가", 7.5, MUTED, False)])

x2 = CX + chw + 0.3
add_rect(s, x2, chy, chw, chh, WHITE, line=LINE)
add_text(s, x2 + 0.16, chy + 0.08, chw - 0.32, 0.22,
         [("전세사기 피해자 누적 결정 건수", 10, INK, True)])
cd2 = CategoryChartData()
cd2.categories = ["’25.8", "’25.10", "’25.11", "’26.5"]
cd2.add_series("누적 결정(건)", (32185, 33978, 34481, 39121))
gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                         IN(x2 + 0.14), IN(chy + 0.34), IN(chw - 0.28), IN(chh - 0.68), cd2)
chart2 = gf2.chart
chart2.has_legend = False
style_chart(chart2)
ser2 = chart2.plots[0].series[0]
ser2.format.line.color.rgb = BLUE7
ser2.format.line.width = Pt(2.25)
chart2.plots[0].has_data_labels = True
chart2.plots[0].data_labels.font.size = Pt(8)
chart2.plots[0].data_labels.font.color.rgb = INK
chart2.plots[0].data_labels.font.name = FONT
chart2.value_axis.visible = False
chart2.value_axis.has_major_gridlines = False
chart2.category_axis.format.line.color.rgb = LINE
add_text(s, x2 + 0.16, chy + chh - 0.28, chw - 0.32, 0.2,
         [("출처: 국토교통부 전세사기피해지원위원회 결정 현황 (보도자료 ’25.8 / ’25.10 / ’25.11 / ’26.5)", 7.5, MUTED, False)])

by = chy + chh + 0.14
add_rect(s, CX, by, CW, 0.96, DARK, radius=0.10)
cols = [
    ("누가 겪는 문제인가", "지역 소상공인 · 전세 계약을 앞둔 청년 · 고령 고객,\n그리고 이들을 담당하는 JB RM · 심사 · 준법 담당자"),
    ("무엇이 문제인가", "기사 · 사기 경보 · 시세 · 등기 · 상담기록이 흩어져 있어\n위험 신호를 조기에 모아서 판단하고 행동하기 어려움"),
    ("해결 후 기대 변화 (KPI)", "위험 인지→대응 착수 시간 50% 단축 · 판단 100% 근거 연결\n고객 대상 행동 100% 사람 승인 통과 · 사후관리 누락 0건"),
]
cw3 = (CW - 0.9) / 3
for i, (t, b) in enumerate(cols):
    x = CX + 0.25 + i * (cw3 + 0.2)
    add_text(s, x, by + 0.12, cw3, 0.24, [(t, 10, CYAN, True)])
    add_text(s, x, by + 0.38, cw3, 0.52, [(line, 8.5, ON_DARK_SOFT, False) for line in b.split("\n")], space_after=1)

# ================================================================ Slide 4 — 솔루션 개요
s = S[3]
remove_guide_boxes(s)

add_text(s, CX, CY - 0.02, CW, 0.24,
         [[("시스템 아키텍처  ", 11, INK, True),
           ("— 케이스 하나가 흐르는 운영 계약: 판단 · 행동 · 검증을 거쳐 감사 기록까지", 9, MUTED, False)]])

# (1) full-width native system pipeline
pipe = [
    ("케이스", "고객 위험", TINT, INK, MUTED, LINE),
    ("에이전트 실행", "AgentRun 로그", TINT, INK, MUTED, LINE),
    ("전문 에이전트", "14종 메시", TINT, INK, MUTED, LINE),
    ("스킬 장착", "25종 저장소", TINT, INK, MUTED, LINE),
    ("근거", "출처 · Evidence", TINT, INK, MUTED, LINE),
    ("PII 거버넌스", "비반출 게이트", BLUE, WHITE, ON_DARK, None),
    ("승인 게이트", "사람 승인", BLUE, WHITE, ON_DARK, None),
    ("감사 원장", "무결성 해시", DARK, ON_DARK, ON_DARK_SOFT, None),
]
n = len(pipe)
pgap = 0.12
pnw = (CW - pgap * (n - 1)) / n
py = CY + 0.30
pnh = 0.66
for i, (t, sub, fill, tc, sc, ln) in enumerate(pipe):
    x = CX + i * (pnw + pgap)
    add_node(s, x, py, pnw, pnh, t, sub, fill, tc, sc, line=ln, title_size=8.5, sub_size=6.5)
    if i < n - 1:
        add_arrow(s, x + pnw, py + pnh / 2, x + pnw + pgap, py + pnh / 2, BLUE, 1.2)
add_text(s, CX, py + pnh + 0.04, CW, 0.18,
         [[("승인 거절 시 ", 7.5, MUTED, False), ("에스컬레이션(상위 검토·외부 차단)", 7.5, BLUE7, True),
           ("으로 분기되며, 모든 단계가 감사 원장에 기록됩니다.", 7.5, MUTED, False)]])

# (2) left — 운영 루프 4단계
sec_y = py + pnh + 0.30
lw = 5.55
add_text(s, CX, sec_y, lw, 0.22, [("운영 루프 — 케이스가 흐르는 4단계", 10.5, INK, True)])
loop_steps = [
    ("1", "케이스 생성", "고객 위험을 케이스로 등록 (RM 또는 정기 스캔)"),
    ("2", "에이전트 실행", "스킬을 장착하고 근거 수집 · 위험 판단"),
    ("3", "사람 승인", "조치 초안을 사람이 검토 — 승인 / 반려"),
    ("4", "감사 기록", "모든 판단 · 행동 · 승인이 활동 이력으로 남음"),
]
ly = sec_y + 0.28
for num, t, b in loop_steps:
    add_rect(s, CX, ly, lw, 0.42, SOFT, line=LINE)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, IN(CX + 0.10), IN(ly + 0.095), IN(0.23), IN(0.23))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    circ.shadow.inherit = False
    tfc = circ.text_frame
    tfc.margin_left = tfc.margin_right = tfc.margin_top = tfc.margin_bottom = 0
    pc = tfc.paragraphs[0]
    pc.alignment = PP_ALIGN.CENTER
    set_run(pc.add_run(), num, 9, WHITE, True)
    add_text(s, CX + 0.46, ly + 0.02, 1.25, 0.38, [(t, 9.5, INK, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, CX + 1.72, ly + 0.02, lw - 1.85, 0.38, [(b, 8, MUTED, False)], anchor=MSO_ANCHOR.MIDDLE)
    ly += 0.48

# (3) right — 에이전트 조직 (사람 승인자 + AI 에이전트)
rx = CX + lw + 0.30
rw = CW - lw - 0.30
add_text(s, rx, sec_y, rw, 0.22,
         [[("에이전트 조직  ", 10.5, INK, True), ("— 사람 승인자 2 + AI 에이전트 14 · 스킬 25", 8.5, MUTED, False)]])
oy = sec_y + 0.28
ha_w = (rw - 0.14) / 2
add_node(s, rx, oy, ha_w, 0.36, "RM 최종 승인자", "사람", NAVY8, WHITE, ON_DARK_SOFT, title_size=8.5, sub_size=6.5)
add_node(s, rx + ha_w + 0.14, oy, ha_w, 0.36, "준법 최종 승인자", "사람", NAVY8, WHITE, ON_DARK_SOFT, title_size=8.5, sub_size=6.5)
oy2 = oy + 0.44
add_node(s, rx, oy2, rw, 0.34, "운영 조율 에이전트 (Orchestrator)", "", DARK, ON_DARK, ON_DARK_SOFT, title_size=8.5)
oy3 = oy2 + 0.42
agents6 = ["위험신호 조기감지", "상환위험 분류", "정책금융 매칭", "RM 보좌", "이상거래 탐지·차단", "준법 검토"]
acw = (rw - 0.20) / 3
for i, name in enumerate(agents6):
    col, row = i % 3, i // 3
    x = rx + col * (acw + 0.10)
    y = oy3 + row * 0.34
    add_chip(s, x, y, acw, 0.28, name, TINT, INK, size=7.5, bold=True, line=LINE)
oy4 = oy3 + 0.74
add_rect(s, rx, oy4, rw, 0.46, CYAN1, line=LINE)
add_text(s, rx + 0.12, oy4 + 0.05, rw - 0.24, 0.18,
         [[("전세 보호 라인  ", 8, BLUE7, True), ("— 전세위험 관리 리드 + 4개 전문 에이전트", 7.5, MUTED, False)]])
add_text(s, rx + 0.12, oy4 + 0.24, rw - 0.24, 0.18,
         [("전세가율 · 등기 권리 · 임차인 손실위험 · 계약 체크리스트 · 은행 연계", 7.5, BODY, False)])

# (4) bottom — design principles strip (dark)
strip_y = CY + CH - 0.30
add_rect(s, CX, strip_y, CW, 0.30, DARK, radius=0.5)
add_text(s, CX + 0.3, strip_y + 0.055, CW - 0.6, 0.22,
         [[("설계 원칙  ", 8.5, CYAN, True),
           ("스킬 장착형 멀티 에이전트(기능 추가 = 스킬 추가) · 승인 우선 자동화(완전 자동 발송 금지) · 근거와 감사 기록은 일급 객체",
            8.5, ON_DARK, False)]])

# ================================================================ Slide 5 — 주요 기능
s = S[4]
remove_guide_boxes(s)

shot_w = 4.70
shot_h = shot_w * 9 / 16
add_text(s, CX, CY - 0.02, shot_w + 0.3, 0.24, [("실제 동작하는 MVP 콘솔", 11, INK, True)])
add_picture_card(s, "shots/dashboard.png", CX, CY + 0.28, shot_w, shot_h)
add_text(s, CX, CY + 0.32 + shot_h, shot_w, 0.2,
         [("대시보드 — 한 줄 지시 실행 · 실시간 실행 · 운영 비용 해석 · 우측 케이스 상세", 7.5, MUTED, False)])
mini_y = CY + 0.56 + shot_h
mini_w = (shot_w - 0.14) / 2
mini_h = mini_w * 9 / 16
add_picture_card(s, "shots/orgchart.png", CX, mini_y, mini_w, mini_h)
add_picture_card(s, "shots/plugins.png", CX + mini_w + 0.14, mini_y, mini_w, mini_h)
add_text(s, CX, mini_y + mini_h + 0.03, mini_w, 0.2, [("에이전트 조직도 (14종)", 7.5, MUTED, False)])
add_text(s, CX + mini_w + 0.14, mini_y + mini_h + 0.03, mini_w, 0.2, [("플러그인 · MCP 커넥터 (6종)", 7.5, MUTED, False)])

fx = CX + shot_w + 0.35
fw = CW - shot_w - 0.35
add_text(s, fx, CY - 0.02, fw, 0.24, [("핵심 기능 6가지 — 모두 MVP에서 동작 확인", 11, INK, True)])
features = [
    ("케이스 보드", "위험 케이스를 칸반/큐로 관리, 클릭 시 근거·담당 에이전트 연동", "흩어진 고객 위험의 단일 작업 화면"),
    ("실행 이력(에이전트 로그)", "실행/지시마다 로그 생성, 상태 전이 자동 기록", "AI 판단 과정의 투명한 추적"),
    ("승인 (사람 검토)", "승인/반려 클릭으로 상태·감사·큐 숫자 즉시 갱신", "고객 대상 행동의 사람 통제"),
    ("전세 보호", "전세가율·권리관계·자산 리스크·체크리스트·은행 연계 5대 진단", "전세사기 계약 전 예방"),
    ("데이터 거버넌스", "PII 토큰화·모델 라우팅·반출 스캔을 패널로 시각화", "외부 LLM을 쓰되 원본 PII 비반출"),
    ("활동 이력(감사)", "근거→판단→행동→승인 전 과정을 시간순 기록", "금융권 감사·설명가능성 충족"),
]
fcw = (fw - 0.16) / 2
fch = 1.18
for i, (t, d, p) in enumerate(features):
    col, row = i % 2, i // 2
    x = fx + col * (fcw + 0.16)
    y = CY + 0.28 + row * (fch + 0.12)
    add_rect(s, x, y, fcw, fch, SOFT, line=LINE)
    add_chip(s, x + fcw - 0.78, y + 0.08, 0.68, 0.22, "구현 완료", SUCCESS_BG, SUCCESS, size=7)
    add_text(s, x + 0.12, y + 0.07, fcw - 0.95, 0.24, [(t, 9.5, INK, True)])
    add_text(s, x + 0.12, y + 0.34, fcw - 0.24, 0.46, [(d, 8, MUTED, False)])
    add_text(s, x + 0.12, y + 0.88, fcw - 0.24, 0.24, [[("해결: ", 7.5, BLUE7, True), (p, 7.5, MUTED, False)]])

strip_y = CY + 0.28 + 3 * fch + 2 * 0.12 + 0.12
add_rect(s, fx, strip_y, fw, 0.42, DARK, radius=0.5)
add_text(s, fx + 0.2, strip_y + 0.10, fw - 0.4, 0.22,
         [[("MVP 범위  ", 8.5, CYAN, True),
           ("에이전트 14 · 스킬 25 · 플러그인 6 · 화면 15 · E2E 19 — github.com/LSB-afk/JB-Fin-AI-Challenge",
            8.5, ON_DARK, False)]])

# ================================================================ Slide 6 — 데이터 및 기술 + 거버넌스
s = S[5]
remove_guide_boxes(s)

add_text(s, CX, CY - 0.02, CW, 0.24,
         [[("데이터 및 기술  ", 11, INK, True),
           ("— 외부 LLM을 쓰되 고객 원본 PII는 외부로 나가지 않는다 (최대 차별점)", 9, MUTED, False)]])

gov = [
    ("Case 필드", "고객 PII", TINT, INK, MUTED, LINE),
    ("① 데이터 등급제", "public~PII 등급", TINT, INK, MUTED, LINE),
    ("② 토큰화·볼트", "국내 PII 보관", BLUE, WHITE, ON_DARK, None),
    ("③ 모델 라우팅", "국내/외부 분리", BLUE, WHITE, ON_DARK, None),
    ("④ 반출 스캔", "패턴+토큰 검증", BLUE, WHITE, ON_DARK, None),
    ("감사 원장", "무결성 해시", DARK, ON_DARK, ON_DARK_SOFT, None),
]
gn = len(gov)
ggap = 0.16
gnw = (CW - ggap * (gn - 1)) / gn
gy = CY + 0.30
gnh = 0.64
for i, (t, sub, fill, tc, sc, ln) in enumerate(gov):
    x = CX + i * (gnw + ggap)
    add_node(s, x, gy, gnw, gnh, t, sub, fill, tc, sc, line=ln, title_size=8.5, sub_size=6.5)
    if i < gn - 1:
        add_arrow(s, x + gnw, gy + gnh / 2, x + gnw + ggap, gy + gnh / 2, BLUE, 1.2)
add_text(s, CX, gy + gnh + 0.04, CW, 0.18,
         [[("법적 근거  ", 7.5, BLUE7, True),
           ("신용정보법 §40-2(분리보관·재식별 금지) · 개인정보보호법 §28-4/§28-5 · 전자금융감독규정 §15(망분리) · 금융위 망분리 개선 로드맵(2024.8)",
            7.5, MUTED, False)]])

mid_y = gy + gnh + 0.34
dw = 6.30
add_text(s, CX, mid_y, dw, 0.22, [("활용 데이터 — 공개 · 공식 출처 우선, 출처 전부 명시", 10.5, INK, True)])
data_rows = [
    ("공공 · 공식", "HUG 안심전세(시세·전세가율·보증) · 국토부 실거래가 · 금융위 보이스피싱 경보 · 한국은행 ECOS", "공개 자료 · 출처 표기"),
    ("시장 · 언론", "연합뉴스 · 쿠키뉴스 등 지역 경기/금리 부담 기사 (위험 신호 수집 대상)", "링크 단위 근거 연결"),
    ("JB 공식", "JB금융그룹 계열사 현황 · JB-네이버클라우드 AI 업무협약(2025.12)", "사업 연계성 근거"),
    ("내부 운영(가상)", "상담 노트 · 케이스 이력 · 승인 기록 — MVP는 가상/시뮬레이션 데이터", "개인정보 마스킹 전제"),
]
dy = mid_y + 0.28
for t, b, m in data_rows:
    add_rect(s, CX, dy, dw, 0.56, WHITE, line=LINE)
    add_text(s, CX + 0.13, dy + 0.05, 1.20, 0.46, [(t, 8.5, BLUE7, True)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, CX + 1.40, dy + 0.05, dw - 2.70, 0.48, [(b, 7.5, INK, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, CX + dw - 1.28, dy + 0.05, 1.18, 0.48, [(m, 6.5, MUTED, False)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    dy += 0.62

tx = CX + dw + 0.30
tw = CW - dw - 0.30
add_text(s, tx, mid_y, tw, 0.22, [("기술 구성 — 현재 MVP와 본선 목표", 10.5, INK, True)])
card_y = mid_y + 0.28
card_h = 2.48
add_rect(s, tx, card_y, tw, card_h, DARK, radius=0.06)
add_text(s, tx + 0.16, card_y + 0.12, tw - 0.32, 0.2, [("현재 MVP  (예선 · 구현 완료)", 8.5, CYAN, True)])
mvp_lines = [
    "Vanilla JS/CSS 정적 콘솔 — 외부 의존성 0, 브라우저만으로 시연",
    "상태 모델: 케이스 · 실행 · 승인 · 감사 상호 연동",
    "GitHub 공개 · Playwright E2E 19종 통과",
]
my = card_y + 0.36
for line in mvp_lines:
    add_text(s, tx + 0.26, my, tw - 0.40, 0.2, [("· " + line, 7.5, ON_DARK_SOFT, False)])
    my += 0.215
add_text(s, tx + 0.16, my + 0.06, tw - 0.32, 0.2, [("본선 목표 아키텍처", 8.5, CYAN, True)])
tech_lines = [
    "LLM — 위험 판단 요약 · 조치 초안 (상용 API, 토큰만 입력)",
    "RAG — HUG·국토부·금융위 공개자료로 근거 강제 연결",
    "Rule Engine — 승인 레벨(L0–L4) · 금지 행동 판정",
    "Multi-Agent — 오케스트레이터가 전문 에이전트에 스킬 주입",
]
ty2 = my + 0.30
for line in tech_lines:
    add_text(s, tx + 0.26, ty2, tw - 0.40, 0.2, [("· " + line, 7.5, ON_DARK_SOFT, False)])
    ty2 += 0.215

# ================================================================ Slide 7 — 사용자 시나리오
s = S[6]
remove_guide_boxes(s)

add_text(s, CX, CY - 0.02, CW, 0.24,
         [[("유스케이스 상세 — 전세 보호 라인  ", 11, INK, True),
           ("(청년 고객이 전세계약 전에 은행을 통해 위험을 진단받는 흐름)", 9, MUTED, False)]])

actors = [("고객", TINT, INK), ("RM(은행)", NAVY8, WHITE), ("AI 에이전트", DARK, ON_DARK), ("검증·승인", BLUE, WHITE)]
lx = CX
for label, fillc, tc in actors:
    add_chip(s, lx, CY + 0.26, 0.98, 0.24, label, fillc, tc, size=7.5,
             line=LINE if fillc == TINT else None)
    lx += 1.08

steps = [
    ("1", "전세계약 전\n상담 신청", TINT, INK, LINE),
    ("2", "케이스 생성\nJBG-201", NAVY8, WHITE, None),
    ("3", "전세가율·권리관계\n·자산 리스크 진단", DARK, ON_DARK, None),
    ("4", "체크리스트 +\n특약 초안 생성", DARK, ON_DARK, None),
    ("5", "준법 검토\n표현·법률 리스크", BLUE, WHITE, None),
    ("6", "RM 승인\n(사람 결정)", NAVY8, WHITE, None),
    ("7", "안전 계약 가이드\n+ 대출·보증 상담", TINT, INK, LINE),
]
fy = CY + 0.60
fh = 0.92
gap = 0.22
fw7 = (CW - gap * 6) / 7
for i, (num, label, fillc, tc, ln) in enumerate(steps):
    x = CX + i * (fw7 + gap)
    add_rect(s, x, fy, fw7, fh, fillc, line=ln, radius=0.12)
    add_text(s, x + 0.08, fy + 0.06, fw7 - 0.16, 0.2, [(num, 10, tc, True)])
    add_text(s, x + 0.08, fy + 0.28, fw7 - 0.16, 0.56,
             [(seg, 8, tc, True) for seg in label.split("\n")], space_after=0)
    if i < 6:
        add_arrow(s, x + fw7 + 0.015, fy + fh / 2, x + fw7 + gap - 0.015, fy + fh / 2, BLUE, 1.1)

add_text(s, CX, fy + fh + 0.05, CW, 0.2,
         [[("안전 정책: ", 8, BLUE7, True),
           ("법률 확정 표현 금지 · 등기/보증은 원문 확인 · 특약은 초안만 · 은행 연계는 고객 동의 + RM 승인 후 진행",
            8, MUTED, False)]])

sy = fy + fh + 0.30
shot_a_w = 4.20
shot_a_h = shot_a_w * 9 / 16
add_picture_card(s, "shots/jeonse.png", CX, sy, shot_a_w, shot_a_h)
add_text(s, CX, sy + shot_a_h + 0.03, shot_a_w, 0.2,
         [("전세 보호 화면 — 5대 진단과 케이스(JBG-201) 연동", 7.5, MUTED, False)])

shot_b_x = CX + shot_a_w + 0.25
shot_b_w = 3.60
shot_b_h = shot_b_w * 9 / 16
add_picture_card(s, "shots/runs.png", shot_b_x, sy + 0.16, shot_b_w, shot_b_h)
add_text(s, shot_b_x, sy + 0.19 + shot_b_h, shot_b_w, 0.2,
         [("실행 이력 — 에이전트 로그와 상태 전이", 7.5, MUTED, False)])

bx = shot_b_x + shot_b_w + 0.25
bw = CX + CW - bx
add_rect(s, bx, sy, bw, shot_a_h + 0.22, SOFT, line=LINE)
add_text(s, bx + 0.14, sy + 0.09, bw - 0.28, 0.4, [("대표 케이스 — 전주 중앙로 카페 자금압박 (JBG-104)", 9, INK, True)])
sub_steps = [
    "매출 둔화 + 금리 부담 신호 감지 (riskScore 88)",
    "상환 위험 분류 → 정책금융·서류 체크리스트 생성",
    "콜백 초안 작성 → 준법 검토 → RM 승인 후 안내",
    "전 과정 감사 기록 — 누락·책임소재 문제 해소",
]
sy2 = sy + 0.42
for i, line in enumerate(sub_steps, 1):
    add_text(s, bx + 0.14, sy2, bw - 0.28, 0.34,
             [[(f"{i}. ", 8, BLUE7, True), (line, 8, INK, False)]], space_after=0)
    sy2 += 0.47

# ================================================================ Slide 8 — 기대 효과
s = S[7]
remove_guide_boxes(s)

impact = [
    ("50%↓", "위험 인지→대응 착수 시간", "케이스 화면 하나로 원인·근거·다음 행동 파악"),
    ("100%", "판단의 근거 연결률", "모든 판단에 출처 링크 또는 내부 이벤트 연결"),
    ("100%", "고객 대상 행동 승인 통과율", "승인 게이트 없이는 외부 안내 불가"),
    ("0건", "사후관리 누락", "후속 태스크 자동 등록 · 사기 외부발송 차단"),
]
iw = (CW - 0.6) / 4
for i, (num, t, b) in enumerate(impact):
    x = CX + i * (iw + 0.2)
    add_rect(s, x, CY, iw, 1.12, TINT, line=LINE)
    add_rect(s, x, CY, iw, 0.07, BLUE, radius=0.5)
    add_text(s, x + 0.13, CY + 0.12, iw - 0.26, 0.32, [(num, 17, BLUE7, True)])
    add_text(s, x + 0.13, CY + 0.47, iw - 0.26, 0.22, [(t, 8.5, INK, True)])
    add_text(s, x + 0.13, CY + 0.71, iw - 0.26, 0.4, [(b, 7.5, MUTED, False)])
add_text(s, CX, CY + 1.16, CW, 0.18,
         [("* MVP 운영 목표 지표 — 본선에서 실데이터 PoC로 측정·검증 예정", 7.5, MUTED, False)])

ry2 = CY + 1.48
rw2 = 6.6
add_text(s, CX, ry2, rw2, 0.22, [("발전 경로 — 예선에서 고객 서비스까지", 11, INK, True)])
road = [
    ("예선 (완료)", "정적 MVP 콘솔\n시나리오 검증 · GitHub 공개", SUCCESS),
    ("본선", "LLM · RAG 연동\n공공데이터 PoC", BLUE),
    ("내부 파일럿", "전북은행 RM 부서\n승인·감사 정책 검증", NAVY8),
    ("고객 서비스화", "뱅킹 앱 연계\n전세 보호 우선 출시", DARK),
]
rcw = (rw2 - 0.45) / 4
for i, (t, b, c) in enumerate(road):
    x = CX + i * (rcw + 0.15)
    add_rect(s, x, ry2 + 0.26, rcw, 1.04, WHITE, line=LINE)
    add_rect(s, x, ry2 + 0.26, rcw, 0.28, c, radius=0.15)
    add_text(s, x + 0.07, ry2 + 0.305, rcw - 0.14, 0.2, [(t, 8, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, x + 0.09, ry2 + 0.62, rcw - 0.18, 0.6,
             [(seg, 7.5, INK, False) for seg in b.split("\n")], space_after=1)
    if i < 3:
        add_arrow(s, x + rcw + 0.005, ry2 + 0.78, x + rcw + 0.145, ry2 + 0.78, BLUE, 1.1)

kx = CX + rw2 + 0.35
kw = CW - rw2 - 0.35
add_text(s, kx, ry2, kw, 0.22, [("운영 리스크 대응", 11, INK, True)])
risks = [
    ("환각 · 오판", "근거 링크 강제 + 확정 표현 금지 + 사람 승인"),
    ("개인정보 · 보안", "PII 비반출 거버넌스 · 망분리 · 접근 권한 분리"),
    ("책임소재 · 설명가능성", "감사 로그로 판단·승인 주체 전 과정 추적"),
    ("저작권 · 라이선스", "공공·공식 출처 우선, 기사 등은 링크 인용"),
]
ky = ry2 + 0.28
for t, b in risks:
    add_text(s, kx, ky, kw, 0.28,
             [[("· " + t + " — ", 8, INK, True), (b, 8, MUTED, False)]], space_after=0)
    ky += 0.27

by8 = ry2 + 1.46
add_rect(s, CX, by8, CW, 0.72, DARK, radius=0.10)
add_text(s, CX + 0.25, by8 + 0.11, CW - 0.5, 0.24,
         [[("계열사 확장:  ", 9, CYAN, True),
           ("전북은행(지역 SME · 전세 보호)  →  광주은행(가계 · 청년)  →  JB우리캐피탈(할부 · 리스 사후관리)",
            9, ON_DARK, False)]])
add_text(s, CX + 0.25, by8 + 0.40, CW - 0.5, 0.22,
         [("JB금융그룹이 공개한 AI 추진 방향(기업여신 상담·심사·사후관리 AI — 네이버클라우드 업무협약)과 같은 축에서 확장",
           8, ON_DARK_SOFT, False)])

prs.save("[LocalGuard] JB금융그룹 Fin AI Challenge MVP제안서.pptx")
print("saved OK")
